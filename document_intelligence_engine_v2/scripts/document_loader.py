"""
Document loader module - comprehensive extraction from virtually any file type.
Merges Version 1 extraction quality (heading detection, section tracking, OCR fallback)
with Version 2 architecture (class-based, FAISS-compatible output).

Supported formats:
    Rich formats:   PDF, DOCX, XLSX, XLS, PPTX, CSV
    Text formats:   TXT, MD, RST, HTML, HTM, XML, JSON, YAML, YML, TOML, INI, CFG
    Code formats:   PY, JS, TS, JAVA, C, CPP, H, CS, GO, RS, RB, PHP, SWIFT, KT, R, SQL
    Data formats:   EPUB, RTF, ODT, ODS, ODP
    Fallback:       Any other text-based file via raw UTF-8 read
"""

import json
import logging
import re
import csv
from pathlib import Path
from typing import Dict, List, Optional

logger = logging.getLogger(__name__)

CHUNK_SIZE = 1000  # characters per chunk (Version 1 proven value)

# ---------------------------------------------------------------------------
# Optional imports — fail gracefully if library not installed
# ---------------------------------------------------------------------------

try:
    import PyPDF2
except ImportError:
    PyPDF2 = None

try:
    import fitz  # PyMuPDF fallback for PDFs
except ImportError:
    fitz = None

try:
    from pdf2image import convert_from_path
    import pytesseract
    OCR_AVAILABLE = True
except ImportError:
    OCR_AVAILABLE = False

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

try:
    from bs4 import BeautifulSoup
except ImportError:
    BeautifulSoup = None

try:
    import ebooklib
    from ebooklib import epub
except ImportError:
    ebooklib = None
    epub = None

try:
    import striprtf.striprtf as striprtf_module
except ImportError:
    striprtf_module = None


# ---------------------------------------------------------------------------
# Chunking helper (from Version 1)
# ---------------------------------------------------------------------------

def _chunk_text(text: str, source_name: str, page_num: int,
                section: Optional[str] = None) -> List[Dict]:
    """Split text into fixed-size chunks, preserving source/page/section metadata."""
    chunks = []
    text = text.strip().replace("\n", " ")
    text = re.sub(r'\s+', ' ', text)
    if not text:
        return chunks
    for i in range(0, len(text), CHUNK_SIZE):
        chunk = text[i:i + CHUNK_SIZE]
        if chunk.strip():
            chunks.append({
                "source": source_name,
                "page": page_num,
                "section": section,
                "text": chunk,
            })
    return chunks


# ---------------------------------------------------------------------------
# Heading detectors (from Version 1)
# ---------------------------------------------------------------------------

def _detect_heading_pdf(line: str) -> Optional[str]:
    line = line.strip()
    if not line:
        return None
    if line.isupper() and len(line) > 3:
        return line
    if line.endswith(":") and len(line) < 80:
        return line
    return None


def _detect_heading_docx(paragraph) -> Optional[str]:
    if "Heading" in paragraph.style.name:
        return paragraph.text.strip() or None
    return None


# ---------------------------------------------------------------------------
# Per-format extractors
# ---------------------------------------------------------------------------

def _extract_pdf(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name

    # Try PyPDF2 first (Version 1 approach)
    if PyPDF2:
        try:
            with open(file_path, 'rb') as f:
                reader = PyPDF2.PdfReader(f)
                for page_num, page in enumerate(reader.pages, start=1):
                    text = page.extract_text() or ""
                    section = None
                    if text.strip():
                        for line in text.split("\n"):
                            h = _detect_heading_pdf(line)
                            if h:
                                section = h
                        chunks.extend(_chunk_text(text, name, page_num, section))
                    elif OCR_AVAILABLE:
                        images = convert_from_path(file_path, first_page=page_num, last_page=page_num)
                        for img in images:
                            ocr_text = pytesseract.image_to_string(img)
                            chunks.extend(_chunk_text(ocr_text, name, page_num, section))
            if chunks:
                return chunks
        except Exception as e:
            logger.warning(f"PyPDF2 failed for {name}: {e}, trying PyMuPDF...")

    # Fallback to PyMuPDF
    if fitz:
        try:
            doc = fitz.open(file_path)
            for page_num in range(len(doc)):
                text = doc[page_num].get_text()
                chunks.extend(_chunk_text(text, name, page_num + 1))
            doc.close()
            return chunks
        except Exception as e:
            logger.error(f"PyMuPDF also failed for {name}: {e}")

    logger.error(f"Could not extract PDF: {name}")
    return chunks


def _extract_docx(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    if not DocxDocument:
        logger.error("python-docx not installed")
        return chunks
    try:
        doc = DocxDocument(file_path)
        section = None
        for i, para in enumerate(doc.paragraphs, start=1):
            h = _detect_heading_docx(para)
            if h:
                section = h
            if para.text.strip():
                chunks.extend(_chunk_text(para.text, name, i, section))
        # Extract tables
        for table_num, table in enumerate(doc.tables, start=1):
            for row_num, row in enumerate(table.rows, start=1):
                row_text = " | ".join(
                    cell.text.strip() for cell in row.cells if cell.text.strip()
                )
                if row_text:
                    chunks.extend(_chunk_text(row_text, f"{name}:Table{table_num}", row_num))
    except Exception as e:
        logger.error(f"DOCX extraction failed for {name}: {e}")
    return chunks


def _extract_xlsx(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    if not openpyxl:
        logger.error("openpyxl not installed")
        return chunks
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            # Group rows in batches of 20 so related data stays together
            row_buffer = []
            chunk_num = 1
            for row in sheet.iter_rows(values_only=True):
                row_text = " | ".join(str(cell) for cell in row if cell is not None)
                if row_text.strip():
                    row_buffer.append(row_text)
                if len(row_buffer) >= 20:
                    combined = " \n ".join(row_buffer)
                    chunks.extend(_chunk_text(combined, name, chunk_num, section=sheet_name))
                    chunk_num += 1
                    row_buffer = []
            # Flush any remaining rows
            if row_buffer:
                combined = " \n ".join(row_buffer)
                chunks.extend(_chunk_text(combined, name, chunk_num, section=sheet_name))
    except Exception as e:
        logger.error(f"XLSX extraction failed for {name}: {e}")
    return chunks


def _extract_csv(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            reader = csv.reader(f)
            for row_num, row in enumerate(reader, start=1):
                row_text = " | ".join(str(cell) for cell in row if cell)
                if row_text.strip():
                    chunks.extend(_chunk_text(row_text, name, row_num))
    except Exception as e:
        logger.error(f"CSV extraction failed for {name}: {e}")
    return chunks


def _extract_pptx(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    if not Presentation:
        logger.error("python-pptx not installed")
        return chunks
    try:
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            if slide_text:
                combined = " ".join(slide_text)
                chunks.extend(_chunk_text(combined, name, slide_num))
    except Exception as e:
        logger.error(f"PPTX extraction failed for {name}: {e}")
    return chunks


def _extract_html(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        if BeautifulSoup:
            soup = BeautifulSoup(content, 'html.parser')
            text = soup.get_text(separator=' ')
        else:
            # Strip tags with regex if BeautifulSoup not available
            text = re.sub(r'<[^>]+>', ' ', content)
        chunks.extend(_chunk_text(text, name, 1))
    except Exception as e:
        logger.error(f"HTML extraction failed for {name}: {e}")
    return chunks


def _extract_epub(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    if not ebooklib:
        logger.error("ebooklib not installed: pip install EbookLib")
        return chunks
    try:
        book = epub.read_epub(str(file_path))
        page_num = 1
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                content = item.get_content().decode('utf-8', errors='ignore')
                if BeautifulSoup:
                    soup = BeautifulSoup(content, 'html.parser')
                    text = soup.get_text(separator=' ')
                else:
                    text = re.sub(r'<[^>]+>', ' ', content)
                chunks.extend(_chunk_text(text, name, page_num))
                page_num += 1
    except Exception as e:
        logger.error(f"EPUB extraction failed for {name}: {e}")
    return chunks


def _extract_rtf(file_path: Path) -> List[Dict]:
    chunks = []
    name = file_path.name
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            content = f.read()
        if striprtf_module:
            text = striprtf_module.rtf_to_text(content)
        else:
            # Rough RTF stripping
            text = re.sub(r'\\[a-z]+\d* ?', ' ', content)
            text = re.sub(r'[{}]', '', text)
        chunks.extend(_chunk_text(text, name, 1))
    except Exception as e:
        logger.error(f"RTF extraction failed for {name}: {e}")
    return chunks


def _extract_text(file_path: Path) -> List[Dict]:
    """Generic extractor for plain text, code, markdown, JSON, YAML, etc."""
    chunks = []
    name = file_path.name
    try:
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            text = f.read()
        chunks.extend(_chunk_text(text, name, 1))
    except Exception as e:
        logger.error(f"Text extraction failed for {name}: {e}")
    return chunks


# ---------------------------------------------------------------------------
# Format routing table
# ---------------------------------------------------------------------------

TEXT_EXTENSIONS = {
    # Plain text
    '.txt', '.md', '.rst', '.log', '.text',
    # Code
    '.py', '.js', '.ts', '.java', '.c', '.cpp', '.h', '.hpp', '.cs',
    '.go', '.rs', '.rb', '.php', '.swift', '.kt', '.r', '.sql', '.sh',
    '.bash', '.zsh', '.ps1', '.bat', '.cmd', '.lua', '.pl', '.scala',
    '.dart', '.vue', '.jsx', '.tsx', '.coffee',
    # Data / config
    '.json', '.yaml', '.yml', '.toml', '.ini', '.cfg', '.conf', '.env',
    '.xml', '.svg',
}

FORMAT_HANDLERS = {
    '.pdf':  _extract_pdf,
    '.docx': _extract_docx,
    '.doc':  _extract_docx,
    '.xlsx': _extract_xlsx,
    '.xls':  _extract_xlsx,
    '.csv':  _extract_csv,
    '.pptx': _extract_pptx,
    '.ppt':  _extract_pptx,
    '.html': _extract_html,
    '.htm':  _extract_html,
    '.epub': _extract_epub,
    '.rtf':  _extract_rtf,
}


def _extract_file(file_path: Path) -> List[Dict]:
    suffix = file_path.suffix.lower()
    if suffix in FORMAT_HANDLERS:
        return FORMAT_HANDLERS[suffix](file_path)
    if suffix in TEXT_EXTENSIONS:
        return _extract_text(file_path)
    # Last resort — try reading as text anyway
    logger.warning(f"Unknown type {suffix} for {file_path.name}, attempting raw text read")
    return _extract_text(file_path)


# ---------------------------------------------------------------------------
# DocumentLoader class (V2 interface — returns doc dicts compatible with
# ChunkingEngine / EmbeddingEngine / VectorStoreManager)
# ---------------------------------------------------------------------------

class DocumentLoader:
    """
    Loads virtually any document type and returns a list of document dicts.
    Each dict is compatible with the rest of the V2 pipeline.
    """

    def __init__(self):
        pass

    def load_document(self, file_path: str) -> Optional[Dict]:
        file_path = Path(file_path)
        if not file_path.exists():
            logger.error(f"File not found: {file_path}")
            return None

        raw_chunks = _extract_file(file_path)
        if not raw_chunks:
            logger.warning(f"No content extracted from {file_path.name}")
            return None

        # Combine all chunks into a single document text
        # (ChunkingEngine will re-chunk; raw_chunks preserve page/section info
        #  which we embed into the text so context isn't lost)
        full_text_parts = []
        for chunk in raw_chunks:
            section = chunk.get("section")
            page = chunk.get("page", 1)
            text = chunk.get("text", "").strip()
            if section:
                full_text_parts.append(f"[Section: {section} | Page: {page}] {text}")
            else:
                full_text_parts.append(f"[Page: {page}] {text}")

        full_text = "\n".join(full_text_parts)

        doc_id = self._generate_doc_id(file_path)

        return {
            "doc_id": doc_id,
            "filename": file_path.name,
            "text": full_text,
            "metadata": {
                "format": file_path.suffix.lower().lstrip('.'),
                "file_size": file_path.stat().st_size,
                "raw_chunk_count": len(raw_chunks),
            },
        }

    def load_directory(self, directory_path: str) -> List[Dict]:
        directory = Path(directory_path)
        documents = []

        if not directory.is_dir():
            logger.error(f"Not a directory: {directory_path}")
            return documents

        for file_path in directory.iterdir():
            if file_path.is_file() and not file_path.name.startswith('.'):
                doc = self.load_document(file_path)
                if doc:
                    documents.append(doc)
                    logger.info(f"Loaded: {file_path.name}")
                else:
                    logger.warning(f"Skipped (no content): {file_path.name}")

        logger.info(f"Loaded {len(documents)} documents from {directory_path}")
        return documents

    @staticmethod
    def _generate_doc_id(file_path: Path) -> str:
        mtime = file_path.stat().st_mtime
        return f"{file_path.stem}_{int(mtime)}"